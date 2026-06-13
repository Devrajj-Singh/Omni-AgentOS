"""Research endpoints - document upload and library management."""
from __future__ import annotations

import io
import os
import uuid
from datetime import datetime

from fastapi import APIRouter, File, HTTPException, UploadFile
from pydantic import BaseModel

from memory.store import memory_store

router = APIRouter()

DOCUMENT_SESSION_PREFIX = "doc:"
MAX_DOCUMENT_SIZE = 20 * 1024 * 1024

SUPPORTED_EXTENSIONS = {
    ".txt", ".md", ".rst", ".log",
    ".py", ".ts", ".tsx", ".js", ".jsx", ".java", ".go", ".rs",
    ".cpp", ".c", ".h", ".cs", ".rb", ".php", ".swift", ".kt",
    ".sh", ".bash", ".zsh", ".yaml", ".yml", ".toml", ".ini",
    ".env", ".cfg", ".conf",
    ".json", ".csv", ".xml", ".html", ".htm",
    ".pdf", ".docx", ".doc", ".pptx", ".ppt", ".xlsx", ".xls",
}

EXTENSION_LABELS = {
    ".pdf": "PDF",
    ".docx": "Word", ".doc": "Word",
    ".pptx": "PowerPoint", ".ppt": "PowerPoint",
    ".xlsx": "Excel", ".xls": "Excel",
    ".md": "Markdown",
    ".py": "Python",
    ".ts": "TypeScript", ".tsx": "TypeScript",
    ".js": "JavaScript", ".jsx": "JavaScript",
    ".json": "JSON",
    ".csv": "CSV",
    ".yaml": "YAML", ".yml": "YAML",
    ".txt": "Text",
    ".html": "HTML", ".htm": "HTML",
}


class DocumentItem(BaseModel):
    id: str
    name: str
    size: int
    file_type: str
    uploaded_at: str
    chunk_count: int


class DocumentListResponse(BaseModel):
    documents: list[DocumentItem]
    total: int


def _extract_text_plain(content: bytes) -> str:
    """Extract text from plain text/code files."""
    try:
        return content.decode("utf-8")
    except UnicodeDecodeError:
        try:
            return content.decode("latin-1")
        except Exception:
            raise ValueError("Cannot decode file as text. File may be binary.")


def _extract_text_pdf(content: bytes) -> str:
    """Extract text from PDF files."""
    try:
        import PyPDF2
    except ImportError:
        raise ValueError("PDF support unavailable. Run: pip install pypdf2")

    try:
        reader = PyPDF2.PdfReader(io.BytesIO(content))
        pages = []
        for i, page in enumerate(reader.pages):
            text = page.extract_text()
            if text and text.strip():
                pages.append(f"[Page {i + 1}]\n{text.strip()}")

        if not pages:
            raise ValueError(
                "PDF appears to be image-based or scanned. "
                "Text extraction requires text-based PDFs."
            )
        return "\n\n".join(pages)
    except ValueError:
        raise
    except Exception as e:
        raise ValueError(f"PDF read error: {str(e)}")


def _extract_text_docx(content: bytes) -> str:
    """Extract text from Word documents (.docx)."""
    try:
        from docx import Document
    except ImportError:
        raise ValueError("Word support unavailable. Run: pip install python-docx")

    try:
        doc = Document(io.BytesIO(content))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(
                    cell.text.strip() for cell in row.cells if cell.text.strip()
                )
                if row_text:
                    paragraphs.append(row_text)
        return "\n\n".join(paragraphs)
    except Exception as e:
        raise ValueError(f"Word document read error: {str(e)}")


def _extract_text_pptx(content: bytes) -> str:
    """Extract text from PowerPoint presentations (.pptx)."""
    try:
        from pptx import Presentation
    except ImportError:
        raise ValueError("PowerPoint support unavailable. Run: pip install python-pptx")

    try:
        prs = Presentation(io.BytesIO(content))
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            if slide_text:
                slides.append(f"[Slide {i}]\n" + "\n".join(slide_text))
        return "\n\n".join(slides)
    except Exception as e:
        raise ValueError(f"PowerPoint read error: {str(e)}")


def _extract_text_xlsx(content: bytes) -> str:
    """Extract text from Excel spreadsheets (.xlsx)."""
    try:
        import openpyxl
    except ImportError:
        raise ValueError("Excel support unavailable. Run: pip install openpyxl")

    try:
        wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
        sheets = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                row_text = " | ".join(str(cell) for cell in row if cell is not None)
                if row_text.strip():
                    rows.append(row_text)
            if rows:
                sheets.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))
        wb.close()
        return "\n\n".join(sheets)
    except Exception as e:
        raise ValueError(f"Excel read error: {str(e)}")


def _extract_text(content: bytes, ext: str) -> str:
    """Route to the correct extractor based on file extension."""
    if ext == ".pdf":
        return _extract_text_pdf(content)
    if ext in {".docx", ".doc"}:
        return _extract_text_docx(content)
    if ext in {".pptx", ".ppt"}:
        return _extract_text_pptx(content)
    if ext in {".xlsx", ".xls"}:
        return _extract_text_xlsx(content)
    return _extract_text_plain(content)


def _chunk_text(text: str, chunk_size: int = 1000, overlap: int = 150) -> list[str]:
    """Split text into overlapping chunks for better retrieval."""
    text = text.strip()
    if not text:
        return []
    if len(text) <= chunk_size:
        return [text]

    chunks = []
    start = 0
    while start < len(text):
        end = min(start + chunk_size, len(text))
        if end < len(text):
            for sep in ["\n\n", "\n", " "]:
                break_at = text.rfind(sep, start, end)
                if break_at > start + chunk_size // 2:
                    end = break_at
                    break
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        if end >= len(text):
            break

        next_start = max(0, end - overlap)
        if next_start <= start:
            next_start = end
        start = next_start

    return chunks


@router.post("/research/upload", response_model=DocumentItem)
async def upload_document(file: UploadFile = File(...)) -> DocumentItem:
    """Upload a document into the RAG pipeline."""
    if not file.filename:
        raise HTTPException(status_code=400, detail="No filename provided.")

    ext = os.path.splitext(file.filename)[1].lower()
    if ext not in SUPPORTED_EXTENSIONS:
        supported = ", ".join(sorted(SUPPORTED_EXTENSIONS))
        raise HTTPException(
            status_code=400,
            detail=f"Unsupported file type '{ext}'. Supported: {supported}",
        )

    content = await file.read()

    if len(content) == 0:
        raise HTTPException(status_code=400, detail="File is empty.")

    if len(content) > MAX_DOCUMENT_SIZE:
        size_mb = len(content) / (1024 * 1024)
        raise HTTPException(
            status_code=400,
            detail=f"File too large ({size_mb:.1f}MB). Maximum size is 20MB.",
        )

    try:
        text = _extract_text(content, ext)
    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e))

    if not text or not text.strip():
        raise HTTPException(
            status_code=400,
            detail="Could not extract any text content from this file.",
        )

    chunks = _chunk_text(text)
    if not chunks:
        raise HTTPException(status_code=400, detail="File produced no extractable content.")

    doc_id = str(uuid.uuid4())
    session_id = f"{DOCUMENT_SESSION_PREFIX}{doc_id}"
    file_type = EXTENSION_LABELS.get(ext, ext.lstrip(".").upper())

    for i, chunk in enumerate(chunks):
        memory_store.save_turn(
            session_id=session_id,
            user_message=f"[Document: {file.filename} | Type: {file_type} | Chunk {i + 1}/{len(chunks)}]",
            assistant_message=chunk,
            workspace_path=file.filename,
        )

    return DocumentItem(
        id=doc_id,
        name=file.filename,
        size=len(content),
        file_type=file_type,
        uploaded_at=datetime.utcnow().isoformat(),
        chunk_count=len(chunks),
    )


@router.get("/research/documents", response_model=DocumentListResponse)
async def list_documents() -> DocumentListResponse:
    """List all uploaded documents."""
    all_memories = memory_store.get_all(limit=1000)
    doc_map: dict[str, DocumentItem] = {}

    for mem in all_memories:
        meta = mem.get("metadata", {})
        session_id = meta.get("session_id", "")
        if not session_id.startswith(DOCUMENT_SESSION_PREFIX):
            continue

        doc_id = session_id.replace(DOCUMENT_SESSION_PREFIX, "")
        filename = meta.get("workspace_path", "Unknown")
        ext = os.path.splitext(filename)[1].lower()
        file_type = EXTENSION_LABELS.get(ext, ext.lstrip(".").upper())

        if doc_id not in doc_map:
            doc_map[doc_id] = DocumentItem(
                id=doc_id,
                name=filename,
                size=0,
                file_type=file_type,
                uploaded_at=meta.get("timestamp", ""),
                chunk_count=0,
            )
        doc_map[doc_id].chunk_count += 1

    documents = sorted(doc_map.values(), key=lambda d: d.uploaded_at, reverse=True)
    return DocumentListResponse(documents=documents, total=len(documents))


@router.delete("/research/documents/{doc_id}")
async def delete_document(doc_id: str) -> dict:
    """Delete all chunks of a document."""
    session_id = f"{DOCUMENT_SESSION_PREFIX}{doc_id}"
    all_memories = memory_store.get_all(limit=1000)
    deleted = 0

    for mem in all_memories:
        meta = mem.get("metadata", {})
        if meta.get("session_id") == session_id:
            if memory_store.delete(mem["id"]):
                deleted += 1

    if deleted == 0:
        raise HTTPException(status_code=404, detail="Document not found.")

    return {"deleted_chunks": deleted, "doc_id": doc_id}
